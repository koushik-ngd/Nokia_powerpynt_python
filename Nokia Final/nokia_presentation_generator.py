#!/usr/bin/env python3
"""
Nokia Failure Analysis PowerPoint Generator
Automatically creates a professional presentation about Nokia's rise and fall
using python-pptx and supporting libraries for charts and visuals.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from io import BytesIO
import numpy as np

# Nokia Brand Colors
NOKIA_BLUE = RGBColor(18, 65, 145)  # Nokia primary blue
NOKIA_LIGHT_BLUE = RGBColor(0, 119, 204)  # Nokia light blue
NOKIA_WHITE = RGBColor(255, 255, 255)
NOKIA_GRAY = RGBColor(102, 102, 102)
NOKIA_DARK_GRAY = RGBColor(51, 51, 51)

class NokiaPresentationGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.setup_slide_master()
        
    def setup_slide_master(self):
        """Configure default slide layout and styling"""
        # Set slide size to 16:9
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        
   # ...existing code...
    def add_title_slide(self):
        """Create stylish title slide with Nokia branding"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add background rectangle with Nokia blue gradient
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, 
            self.prs.slide_width, self.prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = NOKIA_BLUE
        bg_shape.line.fill.background()
        
        # Add Nokia logo image
        logo_path = os.path.join(os.path.dirname(__file__), "nokia_logo.png")
        if os.path.exists(logo_path):
            slide.shapes.add_picture(
                logo_path, Inches(0.5), Inches(0.5), Inches(2), Inches(1)
            )
        else:
            # Fallback: Rectangle placeholder if image not found
            logo_placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), 
                Inches(2), Inches(1)
            )
            logo_placeholder.fill.solid()
            logo_placeholder.fill.fore_color.rgb = NOKIA_WHITE
            logo_text = logo_placeholder.text_frame.text = "NOKIA LOGO"
            logo_placeholder.text_frame.paragraphs[0].font.size = Pt(14)
            logo_placeholder.text_frame.paragraphs[0].font.color.rgb = NOKIA_BLUE
            logo_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
# ...existing code...
        
        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(2), Inches(2.5), Inches(9), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Nokia: Why Did It Fail?"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = NOKIA_WHITE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(2), Inches(4.5), Inches(9), Inches(1.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Industry Analysis & Lessons Learned"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = NOKIA_LIGHT_BLUE
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Team placeholder
        team_box = slide.shapes.add_textbox(
            Inches(8), Inches(6), Inches(4), Inches(1)
        )
        team_frame = team_box.text_frame
        team_frame.text = "Presented by: [Team Name]"
        team_para = team_frame.paragraphs[0]
        team_para.font.size = Pt(16)
        team_para.font.color.rgb = NOKIA_WHITE
        team_para.alignment = PP_ALIGN.RIGHT
        
    def add_agenda_slide(self):
        """Create agenda slide with bullet points and icons"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Agenda"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = NOKIA_BLUE
        
        # Content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "📊 Company Background & History"
        
        agenda_items = [
            "📈 Peak Performance & Market Leadership",
            "📉 The Beginning of Decline", 
            "🌐 Market Trends & Industry Disruption",
            "⚔️ Competitive Analysis",
            "💻 Software & Strategy Failures",
            "🎯 SWOT Analysis",
            "💰 Financial Collapse",
            "📱 Product Gallery",
            "🎓 Legacy & Lessons Learned",
            "🔮 Future Outlook",
            "❓ Q&A"
        ]
        
        for item in agenda_items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = NOKIA_DARK_GRAY
            
    def add_company_background_slide(self):
        """Create company background with timeline infographic"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Company Background & History"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = NOKIA_BLUE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Timeline milestones
        milestones = [
            ("1865", "Founded as paper mill", Inches(1)),
            ("1960s", "Entered electronics", Inches(3.5)),
            ("1980s", "Focus on telecom", Inches(6)),
            ("1990s", "Mobile phone leader", Inches(8.5)),
            ("2000s", "Smartphone era begins", Inches(11))
        ]
        
        # Draw timeline line
        timeline_line = slide.shapes.add_connector(
            1, Inches(1), Inches(3), Inches(11.5), Inches(3)
        )
        timeline_line.line.color.rgb = NOKIA_BLUE
        timeline_line.line.width = Pt(4)
        
        # Add milestone markers
        for year, event, x_pos in milestones:
            # Marker circle
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, x_pos - Inches(0.15), Inches(2.7), 
                Inches(0.3), Inches(0.3)
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = NOKIA_LIGHT_BLUE
            
            # Year label
            year_box = slide.shapes.add_textbox(
                x_pos - Inches(0.5), Inches(2), Inches(1), Inches(0.5)
            )
            year_frame = year_box.text_frame
            year_frame.text = year
            year_para = year_frame.paragraphs[0]
            year_para.font.size = Pt(14)
            year_para.font.bold = True
            year_para.font.color.rgb = NOKIA_BLUE
            year_para.alignment = PP_ALIGN.CENTER
            
            # Event label
            event_box = slide.shapes.add_textbox(
                x_pos - Inches(0.75), Inches(3.5), Inches(1.5), Inches(1)
            )
            event_frame = event_box.text_frame
            event_frame.text = event
            event_para = event_frame.paragraphs[0]
            event_para.font.size = Pt(12)
            event_para.font.color.rgb = NOKIA_DARK_GRAY
            event_para.alignment = PP_ALIGN.CENTER
            
        # Add key facts box
        facts_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(5), 
            Inches(11), Inches(2)
        )
        facts_box.fill.solid()
        facts_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
        facts_box.line.color.rgb = NOKIA_BLUE
        
        facts_text = facts_box.text_frame
        facts_text.text = "Key Facts:\n• Founded in Finland by Fredrik Idestam\n• Transitioned from paper → rubber → electronics → telecom\n• Became world's largest mobile phone vendor (1998-2012)\n• At peak: 40% global market share in mobile phones"
        
        for paragraph in facts_text.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = NOKIA_DARK_GRAY
            
    def create_chart_image(self, chart_type, data, title, filename):
        """Generate chart images using matplotlib"""
        plt.style.use('default')
        fig, ax = plt.subplots(figsize=(10, 6))
        fig.patch.set_facecolor('white')
        
        if chart_type == 'bar':
            bars = ax.bar(data['labels'], data['values'], 
                         color=['#124191', '#0077CC', '#FF6B35', '#4CAF50'])
            ax.set_ylabel('Market Share (%)')
        elif chart_type == 'line':
            ax.plot(data['years'], data['revenue'], 
                   color='#124191', linewidth=3, marker='o')
            ax.set_ylabel('Revenue (Billion EUR)')
            ax.set_xlabel('Year')
        elif chart_type == 'pie':
            colors = ['#124191', '#0077CC', '#FF6B35', '#4CAF50', '#FFC107']
            ax.pie(data['values'], labels=data['labels'], 
                  colors=colors, autopct='%1.1f%%', startangle=90)
            
        ax.set_title(title, fontsize=16, fontweight='bold', color='#124191')
        plt.tight_layout()
        
        # Save to BytesIO
        img_stream = BytesIO()
        plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
        img_stream.seek(0)
        plt.close()
        
        return img_stream
        
    def add_peak_decline_slide(self):
        """Create slide showing Nokia's peak and decline"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Peak Performance & The Beginning of Decline"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = NOKIA_BLUE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Peak section (left side)
        peak_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), 
            Inches(5.5), Inches(5)
        )
        peak_box.fill.solid()
        peak_box.fill.fore_color.rgb = RGBColor(230, 255, 230)
        peak_box.line.color.rgb = RGBColor(76, 175, 80)
        
        peak_title = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.7), Inches(5), Inches(0.5)
        )
        peak_title.text_frame.text = "📈 PEAK (1998-2007)"
        peak_title.text_frame.paragraphs[0].font.size = Pt(20)
        peak_title.text_frame.paragraphs[0].font.bold = True
        peak_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(76, 175, 80)
        
        peak_content = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.3), Inches(5), Inches(3.5)
        )
        peak_text = peak_content.text_frame
        peak_text.text = "• World's largest mobile phone vendor\n• 40% global market share at peak\n• Dominated feature phone era\n• Strong in emerging markets\n• Revenue: €51.1B (2007)\n• 132,000 employees worldwide\n• Symbian OS market leader"
        
        for paragraph in peak_text.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = NOKIA_DARK_GRAY
            
        # Decline section (right side)
        decline_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(7), Inches(1.5), 
            Inches(5.5), Inches(5)
        )
        decline_box.fill.solid()
        decline_box.fill.fore_color.rgb = RGBColor(255, 230, 230)
        decline_box.line.color.rgb = RGBColor(244, 67, 54)
        
        decline_title = slide.shapes.add_textbox(
            Inches(7.2), Inches(1.7), Inches(5), Inches(0.5)
        )
        decline_title.text_frame.text = "📉 DECLINE (2007-2014)"
        decline_title.text_frame.paragraphs[0].font.size = Pt(20)
        decline_title.text_frame.paragraphs[0].font.bold = True
        decline_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(244, 67, 54)
        
        decline_content = slide.shapes.add_textbox(
            Inches(7.2), Inches(2.3), Inches(5), Inches(3.5)
        )
        decline_text = decline_content.text_frame
        decline_text.text = "• iPhone launch (2007) disrupted market\n• Android adoption accelerated\n• Symbian became obsolete\n• Windows Phone partnership failed\n• Market share dropped to 3%\n• Mobile division sold to Microsoft\n• 25,000+ job cuts"
        
        for paragraph in decline_text.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = NOKIA_DARK_GRAY
            
    def add_market_trends_slide(self):
        """Create slide about market trends and industry disruption"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Market Trends & Industry Disruption"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = NOKIA_BLUE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Create smartphone adoption chart
        chart_data = {
            'years': [2007, 2008, 2009, 2010, 2011, 2012, 2013, 2014],
            'revenue': [19, 23, 31, 42, 55, 68, 78, 85]
        }
        
        chart_stream = self.create_chart_image(
            'line', chart_data, 
            'Global Smartphone Adoption (%)', 
            'smartphone_adoption.png'
        )
        
        # Add chart to slide
        chart_pic = slide.shapes.add_picture(
            chart_stream, Inches(1), Inches(1.5), 
            Inches(6), Inches(3.5)
        )
        
        # Add key trends text
        trends_box = slide.shapes.add_textbox(
            Inches(7.5), Inches(1.5), Inches(5), Inches(5)
        )
        trends_text = trends_box.text_frame
        trends_text.text = "Key Industry Disruptions:\n\n🔄 Feature phones → Smartphones\n\n📱 Hardware focus → Software ecosystems\n\n🏪 Carrier control → App stores\n\n🌐 Closed systems → Open platforms\n\n💡 Innovation cycles accelerated\n\n🎯 User experience became critical"
        
        for paragraph in trends_text.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = NOKIA_DARK_GRAY
            
    def add_competitive_analysis_slide(self):
        """Create competitive analysis with charts"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Competitive Analysis: Market Share Evolution"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = NOKIA_BLUE
        title_para.alignment = PP_ALIGN.CENTER
        
        # 2007 Market Share (Nokia's peak)
        data_2007 = {
            'labels': ['Nokia', 'Samsung', 'Apple', 'Others'],
            'values': [40, 15, 0, 45]
        }
        
        chart_2007 = self.create_chart_image(
            'pie', data_2007, '2007 Mobile Phone Market Share', '2007_market.png'
        )
        
        pic_2007 = slide.shapes.add_picture(
            chart_2007, Inches(0.5), Inches(1.5), 
            Inches(5.5), Inches(3.5)
        )
        
        # 2014 Market Share (Nokia's decline)
        data_2014 = {
            'labels': ['Samsung', 'Apple', 'Nokia', 'Huawei', 'Others'],
            'values': [28, 23, 3, 8, 38]
        }
        
        chart_2014 = self.create_chart_image(
            'pie', data_2014, '2014 Smartphone Market Share', '2014_market.png'
        )
        
        pic_2014 = slide.shapes.add_picture(
            chart_2014, Inches(7), Inches(1.5), 
            Inches(5.5), Inches(3.5)
        )
        
        # Add competitor logos placeholders
        competitors = ['Apple', 'Samsung', 'Huawei', 'Ericsson']
        x_positions = [Inches(1), Inches(4), Inches(7), Inches(10)]
        
        for i, (competitor, x_pos) in enumerate(zip(competitors, x_positions)):
            logo_placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x_pos, Inches(5.5), 
                Inches(2), Inches(1)
            )
            logo_placeholder.fill.solid()
            logo_placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
            logo_placeholder.text_frame.text = f"{competitor}\nLOGO"
            logo_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            logo_placeholder.text_frame.paragraphs[0].font.size = Pt(12)
            
    def add_software_strategy_failures_slide(self):
        """Create slide about software and strategy failures"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Software & Strategy Failures"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = NOKIA_BLUE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Create flowchart showing failure cascade
        failure_points = [
            ("Symbian OS", "Outdated architecture", Inches(2), Inches(2)),
            ("App Ecosystem", "Limited developer support", Inches(6), Inches(2)),
            ("User Experience", "Complex interface", Inches(10), Inches(2)),
            ("Innovation Speed", "Slow release cycles", Inches(2), Inches(4)),
            ("Strategic Partnerships", "Microsoft dependency", Inches(6), Inches(4)),
            ("Market Response", "Too little, too late", Inches(10), Inches(4))
        ]
        
        for title, description, x, y in failure_points:
            # Failure box
            failure_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x - Inches(0.75), y, 
                Inches(1.5), Inches(1.2)
            )
            failure_box.fill.solid()
            failure_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
            failure_box.line.color.rgb = RGBColor(244, 67, 54)
            
            # Title text
            title_text = slide.shapes.add_textbox(
                x - Inches(0.7), y + Inches(0.1), 
                Inches(1.4), Inches(0.4)
            )
            title_text.text_frame.text = title
            title_text.text_frame.paragraphs[0].font.size = Pt(12)
            title_text.text_frame.paragraphs[0].font.bold = True
            title_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(244, 67, 54)
            title_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Description text
            desc_text = slide.shapes.add_textbox(
                x - Inches(0.7), y + Inches(0.5), 
                Inches(1.4), Inches(0.6)
            )
            desc_text.text_frame.text = description
            desc_text.text_frame.paragraphs[0].font.size = Pt(10)
            desc_text.text_frame.paragraphs[0].font.color.rgb = NOKIA_DARK_GRAY
            desc_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
        # Add arrows showing cascade effect
        arrow_positions = [
            (Inches(3.5), Inches(2.5), Inches(4.5), Inches(2.5)),
            (Inches(7.5), Inches(2.5), Inches(8.5), Inches(2.5)),
            (Inches(2.75), Inches(3.2), Inches(2.75), Inches(3.8)),
            (Inches(6.75), Inches(3.2), Inches(6.75), Inches(3.8))
        ]
        
        for x1, y1, x2, y2 in arrow_positions:
            arrow = slide.shapes.add_connector(1, x1, y1, x2, y2)
            arrow.line.color.rgb = RGBColor(244, 67, 54)
            arrow.line.width = Pt(2)
            
    def add_swot_analysis_slide(self):
        """Create SWOT analysis with 2x2 table design"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "SWOT Analysis"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = NOKIA_BLUE
        title_para.alignment = PP_ALIGN.CENTER
        
        # SWOT quadrants
        quadrants = [
            ("STRENGTHS", RGBColor(76, 175, 80), Inches(1), Inches(1.5),
             "• Strong brand recognition\n• Global distribution network\n• Hardware expertise\n• Patent portfolio\n• Emerging market presence"),
            ("WEAKNESSES", RGBColor(244, 67, 54), Inches(7), Inches(1.5),
             "• Outdated software platform\n• Slow innovation cycles\n• Poor user experience\n• Limited app ecosystem\n• Rigid corporate culture"),
            ("OPPORTUNITIES", RGBColor(33, 150, 243), Inches(1), Inches(4),
             "• 5G infrastructure market\n• IoT and connected devices\n• Enterprise solutions\n• Network equipment\n• Cloud services"),
            ("THREATS", RGBColor(255, 152, 0), Inches(7), Inches(4),
             "• Smartphone market dominance\n• Platform ecosystems (iOS/Android)\n• Chinese competitors\n• Rapid technology changes\n• Consumer preference shifts")
        ]
        
        for title, color, x, y, content in quadrants:
            # Quadrant box
            quad_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, Inches(5), Inches(2.5)
            )
            quad_box.fill.solid()
            quad_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            quad_box.line.color.rgb = color
            quad_box.line.width = Pt(3)
            
            # Title header
            title_header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, Inches(5), Inches(0.5)
            )
            title_header.fill.solid()
            title_header.fill.fore_color.rgb = color
            title_header.line.fill.background()
            
            # Title text
            title_text = slide.shapes.add_textbox(
                x + Inches(0.1), y + Inches(0.05), 
                Inches(4.8), Inches(0.4)
            )
            title_text.text_frame.text = title
            title_text.text_frame.paragraphs[0].font.size = Pt(18)
            title_text.text_frame.paragraphs[0].font.bold = True
            title_text.text_frame.paragraphs[0].font.color.rgb = NOKIA_WHITE
            title_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Content text
            content_text = slide.shapes.add_textbox(
                x + Inches(0.2), y + Inches(0.7), 
                Inches(4.6), Inches(1.6)
            )
            content_text.text_frame.text = content
            for paragraph in content_text.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = NOKIA_DARK_GRAY
                
    def add_financial_collapse_slide(self):
        """Create financial highlights with charts"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Financial Collapse: Revenue & Profit Decline"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = NOKIA_BLUE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Revenue decline chart
        revenue_data = {
            'years': [2007, 2008, 2009, 2010, 2011, 2012, 2013, 2014],
            'revenue': [51.1, 50.7, 40.9, 42.4, 38.7, 30.2, 12.7, 7.3]
        }
        
        revenue_chart = self.create_chart_image(
            'line', revenue_data, 
            'Nokia Revenue Decline (Billion EUR)', 
            'revenue_decline.png'
        )
        
        # Add revenue chart
        revenue_pic = slide.shapes.add_picture(
            revenue_chart, Inches(0.5), Inches(1.5), 
            Inches(6), Inches(3.5)
        )
        
        # Key financial metrics
        metrics_box = slide.shapes.add_textbox(
            Inches(7), Inches(1.5), Inches(5.5), Inches(5)
        )
        metrics_text = metrics_box.text_frame
        metrics_text.text = "Key Financial Metrics:\n\n📉 Revenue Peak (2007): €51.1B\n📉 Revenue Low (2014): €7.3B\n📉 Decline: -86% over 7 years\n\n💰 Mobile Division Sale: €5.4B\n💼 Job Cuts: 25,000+ employees\n📊 Market Cap Loss: ~€100B\n\n🏭 Restructuring Costs: €2.1B\n📱 R&D Investment Wasted: €15B+"
        
        for paragraph in metrics_text.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = NOKIA_DARK_GRAY
            
    def add_product_gallery_slide(self):
        """Create product gallery with placeholders"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Product Gallery: From Glory to Decline"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = NOKIA_BLUE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Product categories
        products = [
            ("Nokia 3310", "Iconic Feature Phone", Inches(1), Inches(1.5)),
            ("Nokia N95", "Smartphone Pioneer", Inches(4.5), Inches(1.5)),
            ("Nokia Lumia 920", "Windows Phone Era", Inches(8), Inches(1.5)),
            ("Nokia Networks", "5G Infrastructure", Inches(11.5), Inches(1.5)),
            ("Symbian OS", "Mobile Operating System", Inches(2.75), Inches(4)),
            ("HERE Maps", "Navigation Platform", Inches(6.25), Inches(4)),
            ("Nokia Bell Labs", "Research & Innovation", Inches(9.75), Inches(4))
        ]
        
        for product, description, x, y in products:
            # Product placeholder
            product_placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, Inches(2.5), Inches(2)
            )
            product_placeholder.fill.solid()
            product_placeholder.fill.fore_color.rgb = RGBColor(245, 245, 245)
            product_placeholder.line.color.rgb = NOKIA_GRAY
            
            # Product name
            name_text = slide.shapes.add_textbox(
                x + Inches(0.1), y + Inches(0.1), 
                Inches(2.3), Inches(0.5)
            )
            name_text.text_frame.text = product
            name_text.text_frame.paragraphs[0].font.size = Pt(14)
            name_text.text_frame.paragraphs[0].font.bold = True
            name_text.text_frame.paragraphs[0].font.color.rgb = NOKIA_BLUE
            name_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Product image placeholder
            img_placeholder = slide.shapes.add_textbox(
                x + Inches(0.5), y + Inches(0.7), 
                Inches(1.5), Inches(0.8)
            )
            img_placeholder.text_frame.text = "[IMAGE]"
            img_placeholder.text_frame.paragraphs[0].font.size = Pt(12)
            img_placeholder.text_frame.paragraphs[0].font.color.rgb = NOKIA_GRAY
            img_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Description
            desc_text = slide.shapes.add_textbox(
                x + Inches(0.1), y + Inches(1.6), 
                Inches(2.3), Inches(0.3)
            )
            desc_text.text_frame.text = description
            desc_text.text_frame.paragraphs[0].font.size = Pt(10)
            desc_text.text_frame.paragraphs[0].font.color.rgb = NOKIA_DARK_GRAY
            desc_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
    def add_lessons_learned_slide(self):
        """Create lessons learned slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Legacy, Impact & Lessons Learned"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = NOKIA_BLUE
        title_para.alignment = PP_ALIGN.CENTER
        
        # Lessons boxes
        lessons = [
            ("🔄 Adapt or Perish", "Technology disruption requires rapid adaptation. Nokia's slow response to smartphones cost them the market.", 
             Inches(1), Inches(1.5), RGBColor(244, 67, 54)),
            ("💻 Software is King", "Hardware excellence isn't enough. Software platforms and ecosystems determine market success.", 
             Inches(7), Inches(1.5), RGBColor(33, 150, 243)),
            ("👥 User Experience Matters", "Complex interfaces and poor UX drive customers away. Simplicity and intuitive design are crucial.", 
             Inches(1), Inches(4), RGBColor(76, 175, 80)),
            ("🚀 Innovation Speed", "Fast iteration and quick market response are essential in tech. Bureaucracy kills innovation.", 
             Inches(7), Inches(4), RGBColor(255, 152, 0))
        ]
        
        for title, content, x, y, color in lessons:
            # Lesson box
            lesson_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, Inches(5), Inches(2)
            )
            lesson_box.fill.solid()
            lesson_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            lesson_box.line.color.rgb = color
            lesson_box.line.width = Pt(3)
            
            # Title
            title_text = slide.shapes.add_textbox(
                x + Inches(0.2), y + Inches(0.1), 
                Inches(4.6), Inches(0.5)
            )
            title_text.text_frame.text = title
            title_text.text_frame.paragraphs[0].font.size = Pt(18)
            title_text.text_frame.paragraphs[0].font.bold = True
            title_text.text_frame.paragraphs[0].font.color.rgb = color
            
            # Content
            content_text = slide.shapes.add_textbox(
                x + Inches(0.2), y + Inches(0.7), 
                Inches(4.6), Inches(1.1)
            )
            content_text.text_frame.text = content
            content_text.text_frame.paragraphs[0].font.size = Pt(14)
            content_text.text_frame.paragraphs[0].font.color.rgb = NOKIA_DARK_GRAY
            
    def add_future_outlook_slide(self):
        """Create future outlook slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Future Outlook: Nokia's Reinvention"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = NOKIA_BLUE
        
        # Content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "🌐 Network Infrastructure Leader"
        
        outlook_items = [
            "📡 5G technology and equipment",
            "☁️ Cloud and edge computing solutions", 
            "🔬 Bell Labs research and innovation",
            "🏭 Industrial IoT and automation",
            "🛡️ Cybersecurity and network security",
            "🎯 Focus on B2B enterprise solutions",
            "💡 Patent licensing and IP monetization",
            "🤝 Strategic partnerships and acquisitions"
        ]
        
        for item in outlook_items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = NOKIA_DARK_GRAY
            
    def add_thank_you_slide(self):
        """Create thank you/Q&A slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Background
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, 
            self.prs.slide_width, self.prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = NOKIA_BLUE
        bg_shape.line.fill.background()
        
        # Thank you text
        thank_you_box = slide.shapes.add_textbox(
            Inches(2), Inches(2), Inches(9), Inches(2)
        )
        thank_you_frame = thank_you_box.text_frame
        thank_you_frame.text = "Thank You!"
        thank_you_para = thank_you_frame.paragraphs[0]
        thank_you_para.font.size = Pt(48)
        thank_you_para.font.bold = True
        thank_you_para.font.color.rgb = NOKIA_WHITE
        thank_you_para.alignment = PP_ALIGN.CENTER
        
        # Q&A text
        qa_box = slide.shapes.add_textbox(
            Inches(2), Inches(4), Inches(9), Inches(1)
        )
        qa_frame = qa_box.text_frame
        qa_frame.text = "Questions & Discussion"
        qa_para = qa_frame.paragraphs[0]
        qa_para.font.size = Pt(28)
        qa_para.font.color.rgb = NOKIA_LIGHT_BLUE
        qa_para.alignment = PP_ALIGN.CENTER
        
        # Contact info
        contact_box = slide.shapes.add_textbox(
            Inches(2), Inches(5.5), Inches(9), Inches(1.5)
        )
        contact_frame = contact_box.text_frame
        contact_frame.text = "📧 contact@team.com\n🌐 github.com/team/nokia-analysis\n📱 @team_handle"
        contact_para = contact_frame.paragraphs[0]
        contact_para.font.size = Pt(18)
        contact_para.font.color.rgb = NOKIA_WHITE
        contact_para.alignment = PP_ALIGN.CENTER
        
        for paragraph in contact_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = NOKIA_WHITE
            paragraph.alignment = PP_ALIGN.CENTER
            
    def generate_presentation(self):
        """Generate the complete presentation"""
        print("Generating Nokia Failure Analysis Presentation...")
        
        # Add all slides
        self.add_title_slide()
        print("✓ Title slide created")
        
        self.add_agenda_slide()
        print("✓ Agenda slide created")
        
        self.add_company_background_slide()
        print("✓ Company background slide created")
        
        self.add_peak_decline_slide()
        print("✓ Peak and decline slide created")
        
        self.add_market_trends_slide()
        print("✓ Market trends slide created")
        
        self.add_competitive_analysis_slide()
        print("✓ Competitive analysis slide created")
        
        self.add_software_strategy_failures_slide()
        print("✓ Software strategy failures slide created")
        
        self.add_swot_analysis_slide()
        print("✓ SWOT analysis slide created")
        
        self.add_financial_collapse_slide()
        print("✓ Financial collapse slide created")
        
        self.add_product_gallery_slide()
        print("✓ Product gallery slide created")
        
        self.add_lessons_learned_slide()
        print("✓ Lessons learned slide created")
        
        self.add_future_outlook_slide()
        print("✓ Future outlook slide created")
        
        self.add_thank_you_slide()
        print("✓ Thank you slide created")
        
        # Save presentation
        filename = "Nokia_Failure_Analysis_PowerPynt.pptx"
        self.prs.save(filename)
        print(f"\n🎉 Presentation saved as '{filename}'")
        print(f"📊 Total slides: {len(self.prs.slides)}")
        print("\n📋 Presentation includes:")
        print("   • Professional Nokia-branded design")
        print("   • Comprehensive failure analysis")
        print("   • Interactive charts and diagrams")
        print("   • SWOT analysis and lessons learned")
        print("   • Future outlook and Q&A")
        
        return filename

def main():
    """Main function to generate the presentation"""
    try:
        # Create presentation generator
        generator = NokiaPresentationGenerator()
        
        # Generate the presentation
        filename = generator.generate_presentation()
        
        print(f"\n✅ SUCCESS: Nokia presentation generated successfully!")
        print(f"📁 File location: {os.path.abspath(filename)}")
        
    except Exception as e:
        print(f"❌ ERROR: Failed to generate presentation")
        print(f"Error details: {str(e)}")
        raise

if __name__ == "__main__":
    main()